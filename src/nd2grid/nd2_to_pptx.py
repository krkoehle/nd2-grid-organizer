#!/usr/bin/env python3
"""
ND2 to PowerPoint Grid Organizer
---------------------------------
Reads Nikon .nd2 microscopy images and creates PowerPoint slides with
images organized in grids on a black background. Phase contrast / brightfield
images are separated from fluorescent channels, and fluorescent channels
are grouped by color.

Usage:
    python nd2_to_pptx.py <input_folder_or_files> [options]

Examples:
    python nd2_to_pptx.py /path/to/nd2/folder
    python nd2_to_pptx.py /path/to/nd2/folder -o output.pptx --cols 4
    python nd2_to_pptx.py file1.nd2 file2.nd2 file3.nd2
"""

import argparse
import os
import sys
from pathlib import Path
from io import BytesIO

import numpy as np
from PIL import Image
import nd2
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Channel names that indicate phase contrast / brightfield / DIC (non-fluorescent)
PHASE_KEYWORDS = [
    "phase", "brightfield", "bf", "dic", "transmitted", "trans",
    "white", "bright", "daylight", "td", "diascopic",
]


def is_phase_channel(channel_name):
    """Check if a channel name indicates phase contrast or brightfield."""
    name_lower = channel_name.lower().strip()
    return any(kw in name_lower for kw in PHASE_KEYWORDS)


def normalize_image(data, percentile_low=0.5, percentile_high=99.5):
    """Normalize image data to 0-255 uint8 using percentile-based contrast."""
    data = data.astype(np.float64)
    lo = np.percentile(data, percentile_low)
    hi = np.percentile(data, percentile_high)
    if hi <= lo:
        hi = lo + 1
    data = np.clip((data - lo) / (hi - lo) * 255, 0, 255)
    return data.astype(np.uint8)


def channel_to_rgb(data, color):
    """Convert single-channel data to an RGB image using the channel's color.

    Args:
        data: 2D numpy array (normalized 0-255 uint8)
        color: tuple (r, g, b) each 0-255
    Returns:
        PIL Image in RGB mode
    """
    r, g, b = color
    rgb = np.zeros((*data.shape, 3), dtype=np.uint8)
    rgb[..., 0] = (data.astype(np.float64) * r / 255).astype(np.uint8)
    rgb[..., 1] = (data.astype(np.float64) * g / 255).astype(np.uint8)
    rgb[..., 2] = (data.astype(np.float64) * b / 255).astype(np.uint8)
    return Image.fromarray(rgb)


def channel_to_grayscale_rgb(data):
    """Convert single-channel data to grayscale RGB image."""
    rgb = np.stack([data, data, data], axis=-1)
    return Image.fromarray(rgb)


def extract_channels(nd2_path):
    """Extract individual channels from an .nd2 file.

    Returns a list of dicts with keys:
        - 'name': channel name (str)
        - 'color': (r, g, b) tuple
        - 'is_phase': bool
        - 'image': PIL Image (RGB)
        - 'source_file': filename
    """
    channels = []
    try:
        with nd2.ND2File(nd2_path) as f:
            sizes = f.sizes
            n_channels = sizes.get("C", 1)

            # Get channel metadata
            channel_meta = []
            if f.metadata and hasattr(f.metadata, "channels"):
                for ch in f.metadata.channels:
                    c = ch.channel
                    name = c.name
                    color = (c.color.r, c.color.g, c.color.b)
                    channel_meta.append({"name": name, "color": color})

            # Read image data - handle Z-stacks by taking max projection
            data = f.asarray()

            # If there's a Z dimension, do max intensity projection
            if "Z" in sizes:
                z_axis = list(sizes.keys()).index("Z")
                data = np.max(data, axis=z_axis)
                # Recalculate sizes without Z
                remaining_keys = [k for k in sizes.keys() if k != "Z"]
                sizes = {k: v for k, v in sizes.items() if k != "Z"}

            # If there are multiple positions (P), handle each
            # For now, take the first position if multi-position
            if "P" in sizes:
                p_axis = list(sizes.keys()).index("P")
                data = data.take(0, axis=p_axis)
                sizes = {k: v for k, v in sizes.items() if k != "P"}

            # If there's a time dimension, take the last timepoint
            if "T" in sizes:
                t_axis = list(sizes.keys()).index("T")
                data = data.take(-1, axis=t_axis)
                sizes = {k: v for k, v in sizes.items() if k != "T"}

            filename = Path(nd2_path).stem

            if n_channels == 1 or "C" not in sizes:
                # Single channel
                norm = normalize_image(data)
                if channel_meta:
                    name = channel_meta[0]["name"]
                    color = channel_meta[0]["color"]
                    phase = is_phase_channel(name)
                    if phase:
                        img = channel_to_grayscale_rgb(norm)
                    else:
                        img = channel_to_rgb(norm, color)
                else:
                    name = "Channel 0"
                    color = (255, 255, 255)
                    phase = False
                    img = channel_to_grayscale_rgb(norm)

                channels.append({
                    "name": name,
                    "color": color,
                    "is_phase": phase,
                    "image": img,
                    "source_file": filename,
                })
            else:
                # Multi-channel: C should be the first non-spatial axis
                c_axis = list(sizes.keys()).index("C")
                for i in range(n_channels):
                    ch_data = data.take(i, axis=c_axis)
                    norm = normalize_image(ch_data)

                    if i < len(channel_meta):
                        name = channel_meta[i]["name"]
                        color = channel_meta[i]["color"]
                    else:
                        name = f"Channel {i}"
                        color = (255, 255, 255)

                    phase = is_phase_channel(name)
                    if phase:
                        img = channel_to_grayscale_rgb(norm)
                    else:
                        img = channel_to_rgb(norm, color)

                    channels.append({
                        "name": name,
                        "color": color,
                        "is_phase": phase,
                        "image": img,
                        "source_file": filename,
                    })

    except Exception as e:
        print(f"  Warning: Could not read {nd2_path}: {e}")

    return channels


def pil_to_pptx_stream(img, max_size=2048):
    """Convert PIL image to a BytesIO stream for pptx insertion."""
    # Resize if very large to keep pptx manageable
    if max(img.size) > max_size:
        img = img.copy()
        img.thumbnail((max_size, max_size), Image.LANCZOS)
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_grid_slide(prs, images_with_labels, title, cols, slide_width, slide_height):
    """Add a slide with a grid of images on a black background.

    Args:
        prs: Presentation object
        images_with_labels: list of (PIL Image, label_str)
        title: slide title text
        cols: number of columns in the grid
        slide_width: slide width in EMU
        slide_height: slide height in EMU
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    n = len(images_with_labels)
    if n == 0:
        return

    rows = (n + cols - 1) // cols

    # Layout constants (in inches)
    margin = 0.3
    title_height = 0.5
    label_height = 0.25
    spacing = 0.15

    usable_w = Inches(13.333 - 2 * margin)
    usable_h = Inches(7.5 - 2 * margin - title_height)

    cell_w = (usable_w - Emu(Inches(spacing).emu * (cols - 1))) // cols
    cell_h = (usable_h - Emu(Inches(spacing).emu * (rows - 1))) // rows
    img_h = cell_h - Inches(label_height)

    # Add title
    from pptx.util import Emu as EmuUtil
    txBox = slide.shapes.add_textbox(
        Inches(margin), Inches(0.1),
        Inches(13.333 - 2 * margin), Inches(title_height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Place images in grid
    for idx, (img, label) in enumerate(images_with_labels):
        row = idx // cols
        col = idx % cols

        x = Inches(margin) + col * (cell_w + Inches(spacing))
        y = Inches(margin + title_height) + row * (cell_h + Inches(spacing))

        # Calculate image dimensions maintaining aspect ratio
        img_w_max = cell_w
        img_h_max = img_h

        aspect = img.width / img.height
        if aspect > (img_w_max / img_h_max):
            # Width-limited
            final_w = img_w_max
            final_h = int(img_w_max / aspect)
        else:
            # Height-limited
            final_h = img_h_max
            final_w = int(img_h_max * aspect)

        # Center image in cell
        x_offset = (cell_w - final_w) // 2
        y_offset = 0

        stream = pil_to_pptx_stream(img)
        slide.shapes.add_picture(stream, x + x_offset, y + y_offset, final_w, final_h)

        # Add label below image
        label_box = slide.shapes.add_textbox(
            x, y + img_h,
            cell_w, Inches(label_height),
        )
        ltf = label_box.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(11)
        lp.font.color.rgb = RGBColor(255, 255, 255)
        lp.alignment = PP_ALIGN.CENTER


def build_presentation(nd2_paths, output_path, cols=4):
    """Main function: reads .nd2 files, organizes channels, creates PowerPoint.

    Organization strategy:
        - One slide per channel type across all files
        - Phase/brightfield images grouped together
        - Each fluorescent channel gets its own slide
    """
    # Widescreen 16:9
    slide_width = Inches(13.333)
    slide_height = Inches(7.5)

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Collect all channels from all files
    all_channels = []
    for path in sorted(nd2_paths):
        print(f"Reading: {Path(path).name}")
        channels = extract_channels(path)
        all_channels.extend(channels)
        for ch in channels:
            print(f"  -> {ch['name']} ({'phase' if ch['is_phase'] else 'fluorescent'})")

    if not all_channels:
        print("No channels extracted. Check your .nd2 files.")
        sys.exit(1)

    # Group channels
    phase_channels = [ch for ch in all_channels if ch["is_phase"]]
    fluor_channels = [ch for ch in all_channels if not ch["is_phase"]]

    # Group fluorescent by channel name
    fluor_groups = {}
    for ch in fluor_channels:
        key = ch["name"]
        if key not in fluor_groups:
            fluor_groups[key] = []
        fluor_groups[key].append(ch)

    # Create slides
    if phase_channels:
        print(f"\nCreating Phase/Brightfield slide ({len(phase_channels)} images)...")
        images_labels = [(ch["image"], ch["source_file"]) for ch in phase_channels]
        add_grid_slide(prs, images_labels, "Phase Contrast / Brightfield", cols, slide_width, slide_height)

    for ch_name, ch_list in sorted(fluor_groups.items()):
        color = ch_list[0]["color"]
        print(f"Creating {ch_name} slide ({len(ch_list)} images)...")
        images_labels = [(ch["image"], ch["source_file"]) for ch in ch_list]
        add_grid_slide(prs, images_labels, ch_name, cols, slide_width, slide_height)

    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print(f"  {len(prs.slides)} slides, {len(all_channels)} total images")


def main():
    parser = argparse.ArgumentParser(
        description="Organize .nd2 microscopy images into PowerPoint grids.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "input", nargs="+",
        help="Path(s) to .nd2 files or a folder containing .nd2 files",
    )
    parser.add_argument(
        "-o", "--output", default="microscopy_grid.pptx",
        help="Output PowerPoint file path (default: microscopy_grid.pptx)",
    )
    parser.add_argument(
        "-c", "--cols", type=int, default=4,
        help="Number of columns in the image grid (default: 4)",
    )
    args = parser.parse_args()

    # Collect .nd2 file paths
    nd2_paths = []
    for inp in args.input:
        p = Path(inp)
        if p.is_dir():
            nd2_paths.extend(sorted(p.glob("*.nd2")))
        elif p.is_file() and p.suffix.lower() == ".nd2":
            nd2_paths.append(p)
        else:
            print(f"Warning: Skipping {inp} (not an .nd2 file or directory)")

    if not nd2_paths:
        print("Error: No .nd2 files found.")
        sys.exit(1)

    print(f"Found {len(nd2_paths)} .nd2 file(s)\n")
    build_presentation(nd2_paths, args.output, cols=args.cols)


if __name__ == "__main__":
    main()
